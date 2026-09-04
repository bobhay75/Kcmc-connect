from __future__ import annotations

import json
from dataclasses import dataclass, asdict
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

try:
    from strands import Agent, tool
except Exception:  # permits local archive/PPT tests without AWS credentials
    Agent = None
    def tool(fn):
        return fn


@dataclass
class SongMatch:
    title: str
    status: str
    source: str | None = None


def normalize(text: str) -> str:
    return "".join(ch.lower() for ch in text if ch.isalnum())


@tool
def find_song_in_archive(title: str, archive_root: str) -> dict:
    """Find an existing editable PowerPoint song deck by normalized title."""
    target = normalize(title)
    root = Path(archive_root)
    if not root.exists():
        return asdict(SongMatch(title, "missing"))
    for path in root.rglob("*.pptx"):
        if target in normalize(path.stem) or normalize(path.stem) in target:
            return asdict(SongMatch(title, "found", str(path)))
    return asdict(SongMatch(title, "missing"))


@tool
def find_song_in_catalog(title: str, catalog_path: str) -> dict:
    """Search a generated KCMC song catalog before scanning source PowerPoints."""
    path = Path(catalog_path)
    if not path.exists():
        return asdict(SongMatch(title, "missing"))
    data = json.loads(path.read_text(encoding="utf-8"))
    target = normalize(title)
    for segment in data.get("segments", []):
        candidate = normalize(segment.get("title", ""))
        if target == candidate or target in candidate or candidate in target:
            return {
                "title": title,
                "status": "found",
                "source": data.get("source_deck"),
                "start_slide": segment.get("start_slide"),
                "end_slide": segment.get("end_slide"),
                "style": {
                    "font": segment.get("dominant_font"),
                    "font_size_pt": segment.get("dominant_font_size_pt"),
                    "alignment": segment.get("dominant_alignment"),
                },
            }
    return asdict(SongMatch(title, "missing"))


@tool
def build_editable_song_deck(title: str, lyrics: str, output_path: str, service_style: str = "Front Porch") -> str:
    """Create a fully editable PPTX draft using the measured KCMC lyric typography."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]
    blocks = [b.strip() for b in lyrics.split("\n\n") if b.strip()]
    for block in blocks or [title]:
        slide = prs.slides.add_slide(layout)
        box = slide.shapes.add_textbox(Inches(0.65), Inches(0.55), Inches(12.03), Inches(6.35))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = block
        p.font.name = "Arial Narrow"
        p.font.size = Pt(60)
        p.alignment = PP_ALIGN.CENTER
    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    return str(out)


@tool
def quality_check_deck(pptx_path: str) -> dict:
    """Check that the deck opens, has slides, contains editable text, and preserves KCMC typography."""
    path = Path(pptx_path)
    if not path.exists():
        return {"ok": False, "errors": ["file_missing"]}
    prs = Presentation(path)
    errors: list[str] = []
    if len(prs.slides) == 0:
        errors.append("no_slides")
    editable_text = 0
    off_style_runs = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False) or not shape.text.strip():
                continue
            editable_text += 1
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not run.text.strip():
                        continue
                    if run.font.name and run.font.name != "Arial Narrow":
                        off_style_runs += 1
    if editable_text == 0:
        errors.append("no_editable_text")
    if off_style_runs:
        errors.append("off_style_font")
    return {
        "ok": not errors,
        "errors": errors,
        "slides": len(prs.slides),
        "editable_text_shapes": editable_text,
        "off_style_runs": off_style_runs,
    }


@tool
def write_approval_manifest(service_name: str, songs: list[dict], output_path: str) -> str:
    """Create the human-approval gate manifest. Nothing is published automatically."""
    payload = {
        "service": service_name,
        "status": "AWAITING_PASTOR_APPROVAL",
        "songs": songs,
        "autopublish": False,
    }
    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(json.dumps(payload, indent=2), encoding="utf-8")
    return str(out)


SYSTEM_PROMPT = """
You are KCMC Service Agent, a production assistant for worship-service preparation.
Your job is to reduce repetitive production work while preserving human authority.
Always search the generated song catalog first, then the existing archive, before creating anything new.
Prefer editable PPTX assets and preserve the measured KCMC service style.
When an asset is missing, create a draft using the approved service style.
Run quality checks before handoff. Never publish, send, or mark a service final without explicit pastor approval.
Return a concise production report including found assets, created assets, QA failures, and items needing judgment.
""".strip()


def make_agent():
    if Agent is None:
        raise RuntimeError("Install strands-agents before running the autonomous agent")
    return Agent(
        system_prompt=SYSTEM_PROMPT,
        tools=[find_song_in_catalog, find_song_in_archive, build_editable_song_deck, quality_check_deck, write_approval_manifest],
        callback_handler=None,
    )


def run_service_request(request: str):
    agent = make_agent()
    return agent(request)
