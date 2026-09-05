from __future__ import annotations

import argparse
import hashlib
import json
import re
from collections import Counter
from dataclasses import dataclass, asdict
from pathlib import Path
from typing import Iterable

from pptx import Presentation

TITLE_MAX_CHARS = 70
TITLE_MAX_LINES = 3
CCLI_RE = re.compile(r"\bCCLI\b", re.I)

@dataclass
class Segment:
    title: str
    start_slide: int
    end_slide: int
    slide_count: int
    blank_slides: int
    dominant_font: str | None
    dominant_font_size_pt: float | None
    dominant_alignment: str | None
    source_deck: str
    source_sha256: str


def slide_lines(slide) -> list[str]:
    lines: list[str] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        for raw in shape.text.splitlines():
            text = " ".join(raw.split()).strip()
            if text:
                lines.append(text)
    return lines


def normalized_title(lines: list[str]) -> str | None:
    if not lines:
        return None
    filtered = [line for line in lines if not CCLI_RE.search(line)]
    if not filtered:
        return None
    candidate = " ".join(filtered).strip()
    if len(candidate) > TITLE_MAX_CHARS or len(filtered) > TITLE_MAX_LINES:
        return None
    return candidate


def looks_like_title(slide, previous_blank: bool, next_has_text: bool) -> bool:
    lines = slide_lines(slide)
    title = normalized_title(lines)
    if not title or not next_has_text:
        return False
    if len(title.split()) <= 8 and len(title) <= 55:
        return True
    return previous_blank and len(title) <= TITLE_MAX_CHARS


def style_fingerprint(slides: Iterable) -> tuple[str | None, float | None, str | None]:
    fonts: Counter[str] = Counter()
    sizes: Counter[float] = Counter()
    aligns: Counter[str] = Counter()
    for slide in slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.alignment is not None:
                    aligns[str(paragraph.alignment)] += 1
                for run in paragraph.runs:
                    if not run.text.strip():
                        continue
                    if run.font.name:
                        fonts[run.font.name] += 1
                    if run.font.size:
                        sizes[round(run.font.size.pt, 1)] += 1
    return (
        fonts.most_common(1)[0][0] if fonts else None,
        sizes.most_common(1)[0][0] if sizes else None,
        aligns.most_common(1)[0][0] if aligns else None,
    )


def sha256(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def extract_catalog(path: Path) -> dict:
    prs = Presentation(str(path))
    text_flags = [bool(slide_lines(s)) for s in prs.slides]
    candidates: list[tuple[int, str]] = []
    for idx, slide in enumerate(prs.slides):
        prev_blank = idx > 0 and not text_flags[idx - 1]
        next_has_text = idx + 1 < len(prs.slides) and text_flags[idx + 1]
        title = normalized_title(slide_lines(slide))
        terminal_card = idx == len(prs.slides) - 1 and bool(title) and title.lower().startswith("thanks for")
        if terminal_card or looks_like_title(slide, prev_blank, next_has_text):
            if title:
                candidates.append((idx + 1, title))

    cleaned: list[tuple[int, str]] = []
    for item in candidates:
        if cleaned and item[0] - cleaned[-1][0] == 1:
            cleaned[-1] = item
        else:
            cleaned.append(item)

    segments: list[Segment] = []
    digest = sha256(path)
    for pos, (start, title) in enumerate(cleaned):
        end = (cleaned[pos + 1][0] - 1) if pos + 1 < len(cleaned) else len(prs.slides)
        if title.lower().startswith(("thanks for", "welcome home")):
            continue
        subset = [prs.slides[i - 1] for i in range(start, end + 1)]
        blank = sum(1 for slide in subset if not slide_lines(slide))
        font, size, alignment = style_fingerprint(subset)
        segments.append(
            Segment(
                title=title,
                start_slide=start,
                end_slide=end,
                slide_count=end - start + 1,
                blank_slides=blank,
                dominant_font=font,
                dominant_font_size_pt=size,
                dominant_alignment=alignment,
                source_deck=path.name,
                source_sha256=digest,
            )
        )

    return {
        "schema_version": 1,
        "source_deck": path.name,
        "slide_count": len(prs.slides),
        "source_sha256": digest,
        "segments": [asdict(s) for s in segments],
        "copyright_note": "Catalog stores titles, slide ranges, and style metadata only; lyric text is intentionally excluded.",
    }


def main() -> None:
    parser = argparse.ArgumentParser(description="Build a searchable KCMC PowerPoint song catalog without exporting lyrics.")
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--out", type=Path, default=Path("song_catalog.json"))
    args = parser.parse_args()
    catalog = extract_catalog(args.pptx)
    args.out.write_text(json.dumps(catalog, indent=2, ensure_ascii=False), encoding="utf-8")
    print(f"Wrote {len(catalog['segments'])} segments to {args.out}")


if __name__ == "__main__":
    main()
